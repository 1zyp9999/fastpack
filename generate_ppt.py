#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
FastPack PPT 生成脚本
AI 编程马拉松 · 项目成果展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 设置宽屏 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色
    PRIMARY_BLUE = RGBColor(41, 98, 255)
    DARK_BLUE = RGBColor(25, 55, 150)
    ACCENT_ORANGE = RGBColor(255, 107, 53)
    SUCCESS_GREEN = RGBColor(0, 184, 148)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    
    # ==================== 封面页 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 25, 40)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FastPack"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "超快跨平台打包工具"
    p2.font.size = Pt(36)
    p2.font.color.rgb = ACCENT_ORANGE
    p2.alignment = PP_ALIGN.CENTER
    
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1))
    tf3 = info_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "AI 编程马拉松 · 项目成果展示"
    p3.font.size = Pt(24)
    p3.font.color.rgb = RGBColor(200, 200, 200)
    p3.alignment = PP_ALIGN.CENTER
    
    # ==================== 第 1 页：项目速览 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "1. 项目速览（30 秒电梯演讲）"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    col_width = Inches(3.8)
    col_height = Inches(4.5)
    start_y = Inches(1.5)
    
    # 列 1
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), start_y, col_width, col_height)
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_GRAY
    box1.line.color.rgb = PRIMARY_BLUE
    box1.line.width = Pt(2)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "核心问题"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_BLUE
    for text in ["", "• 打包工具速度慢", "• 配置复杂繁琐", "• 平台差异大", "", "• Qt Installer 启动慢", "• 需分别配置 Windows/Linux", "• 缺乏智能识别"]:
        p = tf1.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(8)
    
    # 列 2
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), start_y, col_width, col_height)
    box2.fill.solid()
    box2.fill.fore_color.rgb = LIGHT_GRAY
    box2.line.color.rgb = SUCCESS_GREEN
    box2.line.width = Pt(2)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "解决方案"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = SUCCESS_GREEN
    for text in ["", "FastPack", "", "基于 Rust 构建的", "高性能跨平台打包工具", "", "智能项目识别", "自动化构建流水线", "", "10 倍于传统工具的速度"]:
        p = tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(8)
        if text == "FastPack":
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ORANGE
            p.alignment = PP_ALIGN.CENTER
    
    # 列 3
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), start_y, col_width, col_height)
    box3.fill.solid()
    box3.fill.fore_color.rgb = LIGHT_GRAY
    box3.line.color.rgb = ACCENT_ORANGE
    box3.line.width = Pt(2)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "目标用户/场景"
    p3.font.size = Pt(24)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_ORANGE
    for text in ["", "独立开发者", "快速为 Windows/Linux", "用户生成分发包", "", "小型团队", "统一跨平台发布流程", "减少维护成本", "", "CI/CD 流水线", "自动化构建与打包集成"]:
        p = tf3.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(8)
    
    # ==================== 第 2 页：方法与创新 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "2. 方法与创新"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    table = slide.shapes.add_table(5, 3, Inches(0.3), Inches(1.3), Inches(12.733), Inches(3.2)).table
    headers = ["能力维度", "实现方式", "技术亮点"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    rows = [
        ["智能项目识别", "自动检测 8 种项目类型", "基于文件特征的智能匹配"],
        ["自动化构建", "一键执行全流程", "零配置开箱即用"],
        ["智能压缩", "自适应多线程 zstd", "根据 CPU 动态调整"],
        ["增量打包", "文件变更追踪", "避免重复工作"]
    ]
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            tf = cell.text_frame
            tf.word_wrap = True
            for par in tf.paragraphs:
                par.font.size = Pt(14)
                par.font.color.rgb = RGBColor(40, 40, 40)
    
    tech_y = Inches(4.8)
    tech_title = slide.shapes.add_textbox(Inches(0.5), tech_y, Inches(3), Inches(0.5))
    tf = tech_title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心技术栈"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    tech_items = [
        ("Rust 核心引擎", "零成本抽象，内存安全，极致性能"),
        ("Tauri 框架", "轻量级 GUI，比 Electron 快 3-5 倍"),
        ("zstd 压缩算法", "比 gzip 快 3-5 倍"),
        ("Rayon 并行处理", "充分利用多核 CPU")
    ]
    for i, (tech, desc) in enumerate(tech_items):
        y_pos = tech_y + Inches(0.7) + Inches(0.55 * i)
        item = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12), Inches(0.45))
        tf = item.text_frame
        p = tf.paragraphs[0]
        p.text = tech + ": " + desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(60, 60, 60)
    
    # ==================== 第 3 页：AI 赋能的智能化设计 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "3. AI 赋能的智能化设计"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    flow_y = Inches(1.3)
    box_width = Inches(10)
    steps = [
        ("用户输入项目目录", RGBColor(230, 240, 255), PRIMARY_BLUE),
        ("智能识别引擎：扫描特征、匹配策略、提取版本", RGBColor(230, 255, 240), SUCCESS_GREEN),
        ("自动化构建：选择编译器、并行优化、错误诊断", RGBColor(255, 245, 230), ACCENT_ORANGE),
        ("智能打包：生成安装包、创建菜单、路径自适应", RGBColor(240, 230, 255), RGBColor(120, 80, 180))
    ]
    for i, (text, fill_color, line_color) in enumerate(steps):
        y_pos = flow_y + i * Inches(1.3)
        step = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y_pos, box_width, Inches(1.0))
        step.fill.solid()
        step.fill.fore_color.rgb = fill_color
        step.line.color.rgb = line_color
        step.line.width = Pt(2)
        tf = step.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = line_color
        p.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6), y_pos + Inches(1.1), Inches(0.4), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_ORANGE
            arrow.line.fill.background()
    
    adv_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf_adv = adv_box.text_frame
    tf_adv.word_wrap = True
    p_adv = tf_adv.paragraphs[0]
    p_adv.text = "智能化优势"
    p_adv.font.size = Pt(20)
    p_adv.font.bold = True
    p_adv.font.color.rgb = PRIMARY_BLUE
    advantages = ["", "零学习成本", "无需手动配置，自动识别", "", "智能容错", "构建失败时自动诊断", "", "自适应优化", "根据硬件调整并发策略"]
    for text in advantages[1:]:
        p = tf_adv.add_paragraph()
        p.text = text
        if text in ["零学习成本", "智能容错", "自适应优化"]:
            p.font.size = Pt(16)
            p.font.bold = True
        else:
            p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
    
    # ==================== 第 4 页：成果验证 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "4. 成果验证"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    table = slide.shapes.add_table(5, 5, Inches(0.3), Inches(1.3), Inches(9), Inches(2.5)).table
    headers = ["指标", "FastPack", "Qt Installer", "makepkg", "提升"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    rows = [
        ["打包速度", "基准", "10%", "50%", "10x"],
        ["启动时间", "< 50ms", "500ms+", "N/A", "10x"],
        ["内存占用", "低", "高", "中", "3x"],
        ["压缩速度", "zstd", "gzip", "gzip", "3-5x"]
    ]
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            tf = cell.text_frame
            for par in tf.paragraphs:
                par.font.size = Pt(12)
                par.font.color.rgb = RGBColor(40, 40, 40)
                par.alignment = PP_ALIGN.CENTER
        last_cell = table.cell(row_idx, 4)
        last_cell.fill.solid()
        last_cell.fill.fore_color.rgb = RGBColor(230, 250, 240)
        for par in last_cell.text_frame.paragraphs:
            par.font.bold = True
            par.font.color.rgb = SUCCESS_GREEN
    
    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.6), Inches(1.3), Inches(3.4), Inches(3))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = LIGHT_GRAY
    demo_box.line.color.rgb = PRIMARY_BLUE
    demo_box.line.width = Pt(2)
    tf_demo = demo_box.text_frame
    tf_demo.word_wrap = True
    p_demo = tf_demo.paragraphs[0]
    p_demo.text = "可体验成果"
    p_demo.font.size = Pt(20)
    p_demo.font.bold = True
    p_demo.font.color.rgb = PRIMARY_BLUE
    demo_content = ["", "GitHub:", "github.com/yourusername/fastpack", "", "Linux:", "./build.sh build", "./fastpack", "", "Windows:", ".\\build.bat build"]
    for text in demo_content[1:]:
        p = tf_demo.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(4)
    
    effect_y = Inches(4.5)
    effect_title = slide.shapes.add_textbox(Inches(0.3), effect_y, Inches(3), Inches(0.5))
    tf = effect_title.text_frame
    p = tf.paragraphs[0]
    p.text = "关键效果"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    effects = ["支持 8 种项目类型自动识别", "跨平台统一配置", "一键生成原生安装包"]
    for i, effect in enumerate(effects):
        y_pos = effect_y + Inches(0.6) + Inches(0.45 * i)
        item = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(9), Inches(0.35))
        tf = item.text_frame
        p = tf.paragraphs[0]
        p.text = "* " + effect
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 120, 80)
    
    # ==================== 第 5 页：总结与迭代 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "5. 总结与迭代"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    current_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(6), Inches(3.2))
    current_box.fill.solid()
    current_box.fill.fore_color.rgb = RGBColor(230, 255, 240)
    current_box.line.color.rgb = SUCCESS_GREEN
    current_box.line.width = Pt(2)
    tf_curr = current_box.text_frame
    tf_curr.word_wrap = True
    p_curr = tf_curr.paragraphs[0]
    p_curr.text = "当前版本 (v1.0.0)"
    p_curr.font.size = Pt(20)
    p_curr.font.bold = True
    p_curr.font.color.rgb = SUCCESS_GREEN
    features = ["智能项目类型识别 (8 种)", "自动化构建流水线", "跨平台 GUI (Windows/Linux)", "zstd 高速压缩", "多线程并行处理", "自解压安装包生成", "桌面菜单自动集成"]
    for feature in features:
        p = tf_curr.add_paragraph()
        p.text = "  - " + feature
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(5)
    
    future_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.3), Inches(6.3), Inches(5.2))
    future_box.fill.solid()
    future_box.fill.fore_color.rgb = RGBColor(240, 245, 255)
    future_box.line.color.rgb = PRIMARY_BLUE
    future_box.line.width = Pt(2)
    tf_fut = future_box.text_frame
    tf_fut.word_wrap = True
    p_fut = tf_fut.paragraphs[0]
    p_fut.text = "未来迭代计划"
    p_fut.font.size = Pt(20)
    p_fut.font.bold = True
    p_fut.font.color.rgb = PRIMARY_BLUE
    future_items = ["", "深化智能化:", "  - AI 辅助配置生成", "  - 智能依赖分析", "  - 构建预测优化", "  - 云端构建缓存", "  - 自然语言交互", "", "平台扩展:", "  - macOS 支持", "  - 更多构建系统"]
    for text in future_items[1:]:
        p = tf_fut.add_paragraph()
        p.text = text
        if text in ["深化智能化:", "平台扩展:"]:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
        else:
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_before = Pt(4)
    
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.6))
    tf_footer = footer.text_frame
    p_footer = tf_footer.paragraphs[0]
    p_footer.text = "FastPack - 让跨平台打包变得更快、更简单!"
    p_footer.font.size = Pt(22)
    p_footer.font.bold = True
    p_footer.font.color.rgb = ACCENT_ORANGE
    p_footer.alignment = PP_ALIGN.CENTER
    
    output_path = "/home/jack/Downloads/fastpack/FastPack-presentation.pptx"
    prs.save(output_path)
    print("PPT 已生成：" + output_path)
    return output_path

if __name__ == "__main__":
    create_presentation()
